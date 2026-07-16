# -*- coding: utf-8 -*-
"""RAY local memory brain.

Builds a private, zero-network SQLite search index from user-owned files on
this computer. System folders, application caches, credentials, source-control
internals, and temporary build folders are excluded. No file content is sent
to an API. The background scan keeps image metadata light; explicitly selected
images and scanned PDFs are OCR'd locally and added to the same index.
"""

from __future__ import annotations

import argparse
import csv
import html
import io
import json
import os
import re
import socket
import sqlite3
import subprocess
import sys
import time
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable


DEFAULT_DATA_DIR = (
    Path(os.environ.get("LOCALAPPDATA", Path.home() / "AppData" / "Local"))
    / "RAYEnglish"
    / "brain"
)
DATA_DIR = Path(os.environ.get("RAY_BRAIN_DATA_DIR", DEFAULT_DATA_DIR)).expanduser().resolve()
DB_PATH = DATA_DIR / "ray_brain.db"
STATE_PATH = DATA_DIR / "state.json"
CONFIG_PATH = DATA_DIR / "config.json"
STOP_PATH = DATA_DIR / "stop.request"
REFRESH_PATH = DATA_DIR / "refresh.request"
LOG_PATH = DATA_DIR / "brain.log"
PRIORITY_WRITE_DIR = DATA_DIR / "priority_writes"
LOCK_PORT = 8768

MAX_CONTENT_CHARS = 250_000
MAX_TEXT_FILE_BYTES = 8 * 1024 * 1024
MAX_DOCUMENT_BYTES = 60 * 1024 * 1024
COMMIT_EVERY = 1
STATE_EVERY = 20
WATCH_INTERVAL_SECONDS = 900

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".tsv", ".json", ".jsonl", ".yaml", ".yml",
    ".xml", ".html", ".htm", ".css", ".js", ".ts", ".jsx", ".tsx",
    ".py", ".ps1", ".bat", ".cmd", ".sql", ".ini", ".toml", ".log",
}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".hwp", ".hwpx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
METADATA_EXTENSIONS = {
    *IMAGE_EXTENSIONS,
    ".mp3", ".wav", ".m4a", ".mp4", ".mov", ".avi", ".mkv",
    ".zip", ".7z", ".rar", ".doc", ".ppt", ".xls", ".hwpx",
}
INDEX_EXTENSIONS = TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS | METADATA_EXTENSIONS

EXCLUDED_DIR_NAMES = {
    "$recycle.bin", "$winreagent", "system volume information", "windows", "program files",
    "program files (x86)", "programdata", "recovery", "perflogs", "msocache",
    "config.msi", "documents and settings", "onedrivetemp", "intel", "xboxgames",
    "appdata", "node_modules", ".git", ".svn", ".hg", "__pycache__",
    ".cache", ".claude", ".codex", ".npm", ".nuget", ".gradle", ".idea",
    ".vscode", ".vscode-server", ".cursor", ".windsurf", ".ray_ocr", ".venv", "venv",
    ".m2", ".cargo", ".rustup", "site-packages", "temp", "tmp",
    ".ssh", ".gnupg", ".aws", ".azure", ".kube",
    "_launcher_build", "dist", "build", "archive", "_archive", "legacy_archive",
}
EXCLUDED_FILE_NAMES = {
    "ntuser.dat", "ntuser.dat.log1", "ntuser.dat.log2", "usrclass.dat",
    "credentials", ".env", ".env.local", ".env.production",
}


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def log(message: str) -> None:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    with LOG_PATH.open("a", encoding="utf-8") as handle:
        handle.write(f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] {message}\n")


def write_json_atomic(path: Path, data: dict) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    payload = json.dumps(data, ensure_ascii=False, indent=2)
    last_error: OSError | None = None
    for attempt in range(10):
        try:
            path.write_text(payload, encoding="utf-8")
            return
        except PermissionError as exc:
            last_error = exc
            time.sleep(0.04 * (attempt + 1))
    if last_error:
        raise last_error


def read_state() -> dict:
    for _ in range(4):
        try:
            return json.loads(STATE_PATH.read_text(encoding="utf-8"))
        except Exception:
            time.sleep(0.02)
    return {"status": "not_built", "message": "로컬 기억 뇌가 아직 구축되지 않았습니다."}


def fixed_data_roots() -> list[Path]:
    roots = [Path.home()]
    if os.name == "nt":
        try:
            get_drive_type = __import__("ctypes").windll.kernel32.GetDriveTypeW
            system_drive = Path.home().drive.upper()
            system_root = Path(Path.home().anchor)
            for child in system_root.iterdir():
                if not child.is_dir() or child.name.casefold() == "users" or child.name.startswith("$"):
                    continue
                try:
                    if os.stat(child, follow_symlinks=False).st_file_attributes & 0x400:
                        continue
                except (AttributeError, OSError):
                    pass
                if child.name.casefold() not in EXCLUDED_DIR_NAMES:
                    roots.append(child)
            public = Path.home().parent / "Public"
            if public.is_dir():
                roots.append(public)
            for letter in "DEFGHIJKLMNOPQRSTUVWXYZ":
                root = Path(f"{letter}:\\")
                if root.exists() and get_drive_type(str(root)) == 3 and root.drive.upper() != system_drive:
                    roots.append(root)
        except Exception:
            pass
    unique: list[Path] = []
    seen: set[str] = set()
    for root in roots:
        value = os.path.normcase(str(root.resolve()))
        if value not in seen:
            seen.add(value)
            unique.append(root.resolve())
    return unique


def default_config() -> dict:
    configured_roots = os.environ.get("RAY_BRAIN_ROOTS", "").strip()
    roots = (
        [value for value in configured_roots.split(os.pathsep) if value.strip()]
        if configured_roots
        else [str(p) for p in fixed_data_roots()]
    )
    return {
        "roots": roots,
        "watch_interval_seconds": WATCH_INTERVAL_SECONDS,
        "max_content_chars": MAX_CONTENT_CHARS,
        "local_only": True,
        "excluded_directories": sorted(EXCLUDED_DIR_NAMES),
    }


def load_config() -> dict:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    if not CONFIG_PATH.exists():
        write_json_atomic(CONFIG_PATH, default_config())
    try:
        data = json.loads(CONFIG_PATH.read_text(encoding="utf-8"))
    except Exception:
        data = default_config()
    configured_roots = os.environ.get("RAY_BRAIN_ROOTS", "").strip()
    root_values = (
        [value for value in configured_roots.split(os.pathsep) if value.strip()]
        if configured_roots
        else data.get("roots", [])
    )
    roots = []
    users_root = Path.home().parent.resolve()
    for value in root_values:
        path = Path(os.path.expandvars(os.path.expanduser(str(value))))
        if path.exists():
            resolved = path.resolve()
            if path.name.casefold() in EXCLUDED_DIR_NAMES or resolved == users_root:
                continue
            roots.append(resolved)
    data["roots"] = roots or [Path.home()]
    stored = dict(data)
    stored["roots"] = [str(path) for path in data["roots"]]
    stored["excluded_directories"] = sorted(EXCLUDED_DIR_NAMES)
    write_json_atomic(CONFIG_PATH, stored)
    return data


def connect_db(timeout: float = 30) -> tuple[sqlite3.Connection, bool]:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    database_existed = DB_PATH.exists()
    conn = sqlite3.connect(DB_PATH, timeout=timeout)
    conn.execute(f"PRAGMA busy_timeout={max(1000, int(timeout * 1000))}")
    if database_existed:
        tables = {
            row[0]
            for row in conn.execute(
                "SELECT name FROM sqlite_master WHERE type='table' "
                "AND name IN ('files','files_fts')"
            )
        }
        if "files" in tables:
            return conn, "files_fts" in tables

    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA synchronous=NORMAL")
    conn.execute(
        """
        CREATE TABLE IF NOT EXISTS files (
            path TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            ext TEXT NOT NULL,
            kind TEXT NOT NULL,
            size INTEGER NOT NULL,
            mtime_ns INTEGER NOT NULL,
            content TEXT NOT NULL,
            error TEXT NOT NULL DEFAULT '',
            indexed_at TEXT NOT NULL
        )
        """
    )
    conn.execute("CREATE INDEX IF NOT EXISTS idx_files_ext ON files(ext)")
    fts = True
    try:
        conn.execute(
            "CREATE VIRTUAL TABLE IF NOT EXISTS files_fts USING fts5("
            "path UNINDEXED, name, content, tokenize='unicode61 remove_diacritics 2')"
        )
    except sqlite3.OperationalError:
        fts = False
    conn.commit()
    return conn, fts


def decode_text(raw: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-16", "cp949", "latin-1"):
        try:
            return raw.decode(encoding)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", "replace")


def clean_text(text: str, ext: str = "") -> str:
    text = text.replace("\x00", " ")
    if ext in {".html", ".htm"}:
        text = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", text)
        text = re.sub(r"(?s)<[^>]+>", " ", text)
        text = html.unescape(text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{4,}", "\n\n\n", text)
    return text.strip()[:MAX_CONTENT_CHARS]


def read_plain(path: Path) -> str:
    with path.open("rb") as handle:
        return decode_text(handle.read(MAX_TEXT_FILE_BYTES + 1))


def read_pdf(path: Path) -> tuple[str, str]:
    import fitz

    out: list[str] = []
    document = fitz.open(path)
    try:
        for page in document:
            out.append(page.get_text("text"))
            if sum(len(part) for part in out) >= MAX_CONTENT_CHARS:
                break
    finally:
        document.close()
    text = "\n".join(out).strip()
    if not text:
        return "", "스캔 PDF: 개인정보 보호를 위해 외부 OCR을 사용하지 않음"
    return text, ""


def read_docx(path: Path) -> str:
    import docx

    document = docx.Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
            if sum(len(part) for part in parts) >= MAX_CONTENT_CHARS:
                break
    return "\n".join(parts)


def read_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    for number, slide in enumerate(presentation.slides, 1):
        slide_parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text)
                if text:
                    slide_parts.append(text)
        if slide_parts:
            parts.append(f"[Slide {number}]\n" + "\n".join(slide_parts))
        if sum(len(part) for part in parts) >= MAX_CONTENT_CHARS:
            break
    return "\n".join(parts)


def read_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row_number, row in enumerate(sheet.iter_rows(values_only=True), 1):
                values = [str(value) for value in row if value is not None]
                if values:
                    parts.append("\t".join(values))
                if row_number >= 3000 or sum(len(part) for part in parts) >= MAX_CONTENT_CHARS:
                    break
            if sum(len(part) for part in parts) >= MAX_CONTENT_CHARS:
                break
    finally:
        workbook.close()
    return "\n".join(parts)


def read_hwp(path: Path) -> tuple[str, str]:
    import olefile

    ole = olefile.OleFileIO(path)
    try:
        if ole.exists("PrvText"):
            return ole.openstream("PrvText").read().decode("utf-16-le", "ignore"), ""
        return "", "HWP 미리보기 텍스트가 없어 메타데이터만 기억함"
    finally:
        ole.close()


def read_hwpx(path: Path) -> str:
    parts: list[str] = []
    with zipfile.ZipFile(path) as archive:
        names = [name for name in archive.namelist() if name.lower().endswith(".xml") and "section" in name.lower()]
        for name in sorted(names):
            raw = archive.read(name).decode("utf-8", "replace")
            parts.extend(re.findall(r">([^<>]+)<", raw))
            if sum(len(part) for part in parts) >= MAX_CONTENT_CHARS:
                break
    return "\n".join(html.unescape(part) for part in parts)


def extract_content_raw(
    path: Path, ext: str, size: int, *, allow_ocr: bool = False
) -> tuple[str, str, str]:
    if ext in IMAGE_EXTENSIONS:
        if not allow_ocr:
            return "", "image_metadata", ""
        try:
            import ray_ingest

            return clean_text(ray_ingest.ingest(str(path))), "image_ocr", ""
        except Exception as exc:
            return "", "image_ocr", f"로컬 이미지 OCR 실패: {exc}"
    if ext in METADATA_EXTENSIONS and ext not in DOCUMENT_EXTENSIONS:
        return "", "metadata", ""
    if ext in TEXT_EXTENSIONS:
        if size > MAX_TEXT_FILE_BYTES:
            return "", "text", "텍스트 파일 크기 제한 초과"
        return clean_text(read_plain(path), ext), "text", ""
    if size > MAX_DOCUMENT_BYTES:
        return "", "document", "문서 크기 제한 초과"
    try:
        if ext == ".pdf":
            content, error = read_pdf(path)
            if not content and allow_ocr:
                try:
                    import ray_ingest

                    content = ray_ingest.ingest(str(path))
                    error = "" if content else error
                except Exception as exc:
                    error = f"로컬 PDF OCR 실패: {exc}"
            return clean_text(content), "pdf", error
        if ext == ".docx":
            return clean_text(read_docx(path)), "docx", ""
        if ext == ".pptx":
            return clean_text(read_pptx(path)), "pptx", ""
        if ext == ".xlsx":
            return clean_text(read_xlsx(path)), "xlsx", ""
        if ext == ".hwp":
            content, error = read_hwp(path)
            return clean_text(content), "hwp", error
        if ext == ".hwpx":
            return clean_text(read_hwpx(path)), "hwpx", ""
    except Exception as exc:
        return "", ext.lstrip(".") or "document", str(exc)[:240]
    return "", "metadata", ""


def extract_content(
    path: Path, ext: str, size: int, *, allow_ocr: bool = False
) -> tuple[str, str, str]:
    if ext not in DOCUMENT_EXTENSIONS and not (allow_ocr and ext in IMAGE_EXTENSIONS):
        return extract_content_raw(path, ext, size, allow_ocr=allow_ocr)

    creation_flags = subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0
    env = os.environ.copy()
    env["PYTHONUTF8"] = "1"
    try:
        result = subprocess.run(
            [
                sys.executable,
                str(Path(__file__).resolve()),
                "_extract",
                str(path),
                *( ["--ocr"] if allow_ocr else [] ),
            ],
            cwd=str(Path(__file__).resolve().parent),
            env=env,
            capture_output=True,
            text=True,
            encoding="utf-8",
            errors="replace",
            timeout=420 if allow_ocr else 45,
            creationflags=creation_flags,
        )
    except subprocess.TimeoutExpired:
        limit = 420 if allow_ocr else 45
        return "", ext.lstrip(".") or "document", f"문서 추출 제한시간 {limit}초 초과"
    if result.returncode != 0:
        detail = (result.stderr or result.stdout or f"worker exit {result.returncode}").strip()
        return "", ext.lstrip(".") or "document", f"격리 파서 실패: {detail[-180:]}"
    try:
        payload = json.loads(result.stdout)
        return (
            str(payload.get("content", "")),
            str(payload.get("kind", ext.lstrip(".") or "document")),
            str(payload.get("error", "")),
        )
    except Exception as exc:
        return "", ext.lstrip(".") or "document", f"격리 파서 응답 오류: {exc}"


def excluded_directory(path: Path) -> bool:
    name = path.name.casefold()
    if (
        name in EXCLUDED_DIR_NAMES
        or name.startswith("~$")
        or name.startswith("쓰레기통_정리")
        or "중간산출물" in name
    ):
        return True
    try:
        attrs = os.stat(path).st_file_attributes
        if attrs & 0x2 or attrs & 0x4:
            return True
    except (AttributeError, OSError):
        pass
    return False


def iter_files(roots: Iterable[Path]) -> Iterable[Path]:
    for root in roots:
        for current, dirs, files in os.walk(root, topdown=True, followlinks=False):
            current_path = Path(current)
            dirs[:] = [name for name in dirs if not excluded_directory(current_path / name)]
            for name in files:
                if name.casefold() in EXCLUDED_FILE_NAMES or name.startswith("~$"):
                    continue
                path = current_path / name
                if path.suffix.casefold() in INDEX_EXTENSIONS:
                    yield path


def update_state(state: dict, **changes) -> dict:
    state.update(changes)
    state["updated_at"] = utc_now()
    write_json_atomic(STATE_PATH, state)
    return state


def priority_write_pending() -> bool:
    """Yield the background writer while a user-selected file is being saved."""
    if not PRIORITY_WRITE_DIR.exists():
        return False
    now = time.time()
    pending = False
    for marker in PRIORITY_WRITE_DIR.glob("*.lock"):
        try:
            if now - marker.stat().st_mtime > 120:
                marker.unlink(missing_ok=True)
            else:
                pending = True
        except OSError:
            continue
    return pending


def upsert_file(
    conn: sqlite3.Connection,
    fts: bool,
    path: Path,
    stat: os.stat_result,
    content: str,
    kind: str,
    error: str,
) -> None:
    value = str(path)
    ext = path.suffix.casefold()
    stored_content = content if not fts else content[:2000]
    conn.execute(
        """
        INSERT INTO files(path,name,ext,kind,size,mtime_ns,content,error,indexed_at)
        VALUES(?,?,?,?,?,?,?,?,?)
        ON CONFLICT(path) DO UPDATE SET
          name=excluded.name, ext=excluded.ext, kind=excluded.kind,
          size=excluded.size, mtime_ns=excluded.mtime_ns,
          content=excluded.content, error=excluded.error, indexed_at=excluded.indexed_at
        """,
        (value, path.name, ext, kind, stat.st_size, stat.st_mtime_ns, stored_content, error, utc_now()),
    )
    if fts:
        conn.execute("DELETE FROM files_fts WHERE path=?", (value,))
        conn.execute(
            "INSERT INTO files_fts(path,name,content) VALUES(?,?,?)",
            (value, path.name, content),
        )


def index_once() -> dict:
    config = load_config()
    roots: list[Path] = config["roots"]
    REFRESH_PATH.unlink(missing_ok=True)
    state = {
        "status": "indexing",
        "pid": os.getpid(),
        "roots": [str(path) for path in roots],
        "scanned": 0,
        "indexed": 0,
        "updated": 0,
        "skipped": 0,
        "errors": 0,
        "current": "",
        "started_at": utc_now(),
        "message": "사용자 파일을 로컬 기억으로 정리하고 있습니다.",
        "local_only": True,
    }
    write_json_atomic(STATE_PATH, state)
    log("색인 시작: " + ", ".join(str(path) for path in roots))
    conn, fts = connect_db()
    existing = {
        row[0]: (row[1], row[2])
        for row in conn.execute("SELECT path,size,mtime_ns FROM files")
    }
    seen: set[str] = set()
    try:
        for path in iter_files(roots):
            if STOP_PATH.exists():
                state["status"] = "stopping"
                state["message"] = "중지 요청을 처리했습니다."
                break
            value = str(path)
            seen.add(value)
            state["scanned"] += 1
            state["current"] = value
            try:
                stat = path.stat()
                old = existing.get(value)
                if old and old == (stat.st_size, stat.st_mtime_ns):
                    state["skipped"] += 1
                else:
                    if path.suffix.casefold() in DOCUMENT_EXTENSIONS:
                        conn.commit()
                    content, kind, error = extract_content(path, path.suffix.casefold(), stat.st_size)
                    while priority_write_pending() and not STOP_PATH.exists():
                        conn.commit()
                        time.sleep(0.05)
                    upsert_file(conn, fts, path, stat, content, kind, error)
                    # Keep the writer lock short so an explicitly selected file can
                    # be OCR'd and remembered while the full-PC scan is still active.
                    conn.commit()
                    state["indexed"] += 1
                    if old:
                        state["updated"] += 1
                    if error:
                        state["errors"] += 1
                if state["scanned"] % COMMIT_EVERY == 0:
                    conn.commit()
            except (OSError, PermissionError) as exc:
                state["errors"] += 1
                log(f"파일 접근 실패: {path} :: {exc}")
            if state["scanned"] % STATE_EVERY == 0:
                update_state(state)

        if state["status"] != "stopping":
            stale = [path for path in existing if path not in seen]
            for value in stale:
                conn.execute("DELETE FROM files WHERE path=?", (value,))
                if fts:
                    conn.execute("DELETE FROM files_fts WHERE path=?", (value,))
            conn.commit()
            total = conn.execute("SELECT COUNT(*) FROM files").fetchone()[0]
            db_mb = round(DB_PATH.stat().st_size / (1024 * 1024), 2) if DB_PATH.exists() else 0
            update_state(
                state,
                status="ready",
                current="",
                total_files=total,
                db_mb=db_mb,
                finished_at=utc_now(),
                message=f"로컬 기억 {total:,}개 준비 완료",
                fts=fts,
            )
            log(f"색인 완료: total={total}, scanned={state['scanned']}, indexed={state['indexed']}")
        else:
            conn.commit()
            update_state(state, status="stopped", current="", finished_at=utc_now())
    except Exception as exc:
        update_state(state, status="error", current="", message=str(exc)[:300], finished_at=utc_now())
        log(f"색인 오류: {exc}")
        raise
    finally:
        conn.close()
    return state


def acquire_watch_lock() -> socket.socket | None:
    lock = socket.socket(socket.AF_INET, socket.SOCK_STREAM)
    try:
        lock.bind(("127.0.0.1", LOCK_PORT))
        lock.listen(1)
        return lock
    except OSError:
        lock.close()
        return None


def watch() -> int:
    lock = acquire_watch_lock()
    if lock is None:
        log("기존 로컬 뇌가 이미 실행 중")
        return 0
    STOP_PATH.unlink(missing_ok=True)
    try:
        while not STOP_PATH.exists():
            state = index_once()
            if state.get("status") in {"error", "stopped"}:
                break
            state = read_state()
            update_state(
                state,
                status="watching",
                pid=os.getpid(),
                message=f"로컬 기억 {state.get('total_files', 0):,}개 · 자동 갱신 대기",
            )
            interval = int(load_config().get("watch_interval_seconds", WATCH_INTERVAL_SECONDS))
            for _ in range(max(30, interval)):
                if STOP_PATH.exists() or REFRESH_PATH.exists():
                    break
                time.sleep(1)
        state = read_state()
        update_state(state, status="stopped", pid=0, message="로컬 기억 뇌가 중지되었습니다.")
        return 0
    finally:
        STOP_PATH.unlink(missing_ok=True)
        REFRESH_PATH.unlink(missing_ok=True)
        lock.close()


def search(query: str, limit: int = 20) -> list[dict]:
    query = query.strip()
    if not query or not DB_PATH.exists():
        return []
    conn, fts = connect_db()
    try:
        if fts:
            terms = re.findall(r"[0-9A-Za-z가-힣_\-]+", query)
            expression = " AND ".join(f'"{term}"*' for term in terms) or f'"{query}"*'
            rows = conn.execute(
                """
                SELECT path, name,
                       snippet(files_fts, 2, '[', ']', ' … ', 18) AS excerpt,
                       bm25(files_fts) AS score
                FROM files_fts
                WHERE files_fts MATCH ?
                ORDER BY score
                LIMIT ?
                """,
                (expression, max(1, min(limit, 100))),
            ).fetchall()
        else:
            needle = f"%{query}%"
            rows = conn.execute(
                "SELECT path,name,substr(content,1,500),0 FROM files "
                "WHERE name LIKE ? OR content LIKE ? LIMIT ?",
                (needle, needle, max(1, min(limit, 100))),
            ).fetchall()
        return [
            {"path": row[0], "name": row[1], "excerpt": row[2] or "", "score": row[3]}
            for row in rows
        ]
    finally:
        conn.close()


def ingest_one(path_value: str) -> dict:
    path = Path(path_value).expanduser().resolve()
    if not path.is_file():
        raise FileNotFoundError(path)
    ext = path.suffix.casefold()
    if ext not in INDEX_EXTENSIONS:
        raise ValueError(f"지원하지 않는 파일 형식: {ext or '(확장자 없음)'}")
    stat = path.stat()
    content, kind, error = extract_content(path, ext, stat.st_size, allow_ocr=True)
    PRIORITY_WRITE_DIR.mkdir(parents=True, exist_ok=True)
    marker = PRIORITY_WRITE_DIR / f"{os.getpid()}_{time.time_ns()}.lock"
    marker.write_text(str(path), encoding="utf-8")
    conn: sqlite3.Connection | None = None
    try:
        time.sleep(0.15)
        conn, fts = connect_db(timeout=10)
        last_error = None
        for attempt in range(8):
            try:
                conn.execute("BEGIN IMMEDIATE")
                upsert_file(conn, fts, path, stat, content, kind, error)
                conn.commit()
                last_error = None
                break
            except sqlite3.OperationalError as exc:
                conn.rollback()
                if "locked" not in str(exc).casefold():
                    raise
                last_error = exc
                time.sleep(0.15 * (attempt + 1))
        if last_error:
            raise last_error
    finally:
        if conn is not None:
            conn.close()
        marker.unlink(missing_ok=True)
    return {
        "ok": not bool(error),
        "path": str(path),
        "kind": kind,
        "characters": len(content),
        "error": error,
        "local_only": True,
    }


def command_status() -> int:
    print(json.dumps(read_state(), ensure_ascii=False, indent=2))
    return 0


def command_refresh() -> int:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    REFRESH_PATH.touch()
    print(json.dumps({"ok": True, "action": "refresh"}, ensure_ascii=False))
    return 0


def command_stop() -> int:
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    STOP_PATH.touch()
    print(json.dumps({"ok": True, "action": "stop"}, ensure_ascii=False))
    return 0


def main() -> int:
    parser = argparse.ArgumentParser(description="RAY local memory brain")
    sub = parser.add_subparsers(dest="command", required=True)
    sub.add_parser("index")
    sub.add_parser("watch")
    sub.add_parser("status")
    sub.add_parser("refresh")
    sub.add_parser("stop")
    extract_parser = sub.add_parser("_extract", help=argparse.SUPPRESS)
    extract_parser.add_argument("path")
    extract_parser.add_argument("--ocr", action="store_true")
    ingest_parser = sub.add_parser("ingest")
    ingest_parser.add_argument("path")
    search_parser = sub.add_parser("search")
    search_parser.add_argument("query")
    search_parser.add_argument("--limit", type=int, default=20)
    search_parser.add_argument("--json", action="store_true")
    args = parser.parse_args()

    if args.command == "index":
        print(json.dumps(index_once(), ensure_ascii=False, indent=2))
        return 0
    if args.command == "watch":
        return watch()
    if args.command == "status":
        return command_status()
    if args.command == "refresh":
        return command_refresh()
    if args.command == "stop":
        return command_stop()
    if args.command == "_extract":
        path = Path(args.path)
        stat = path.stat()
        content, kind, error = extract_content_raw(
            path, path.suffix.casefold(), stat.st_size, allow_ocr=args.ocr
        )
        print(json.dumps({"content": content, "kind": kind, "error": error}, ensure_ascii=False))
        return 0
    if args.command == "ingest":
        print(json.dumps(ingest_one(args.path), ensure_ascii=False, indent=2))
        return 0
    if args.command == "search":
        results = search(args.query, args.limit)
        if args.json:
            print(json.dumps(results, ensure_ascii=False, indent=2))
        else:
            for item in results:
                print(f"{item['name']}\n  {item['path']}\n  {item['excerpt']}\n")
        return 0
    return 2


if __name__ == "__main__":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")
    raise SystemExit(main())
