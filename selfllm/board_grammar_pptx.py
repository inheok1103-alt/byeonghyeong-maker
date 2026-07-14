# -*- coding: utf-8 -*-
"""board_grammar_pptx.py — 어법 문제 JSON → RAY 어법 수업 보드 PPTX
사용: python board_grammar_pptx.py <grammar.json> <out.pptx>
슬라이드: 01 표지 · 02 지문통독 · 03 실전문항 · 04 밑줄5포인트진단 · 05 정답근거 · 06 직독직해 · 07 핵심문법·어휘·과제"""
import sys, io, json, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY=RGBColor(0x1A,0x2A,0x3A); NAVY2=RGBColor(0x23,0x36,0x48); NAVYD=RGBColor(0x0F,0x1B,0x27)
GOLD=RGBColor(0xC5,0xA8,0x69); GOLDL=RGBColor(0xE8,0xD5,0xA5); IVORY=RGBColor(0xF7,0xF5,0xEF)
INK=RGBColor(0x23,0x2A,0x31); GRAY=RGBColor(0x5B,0x65,0x70); BORDER=RGBColor(0xE7,0xE3,0xD8); WHITE=RGBColor(0xFF,0xFF,0xFF)
F_BLUE=RGBColor(0x3E,0x7C,0xB1); F_ORANGE=RGBColor(0xD2,0x7A,0x3D); F_RED=RGBColor(0xB9,0x4A,0x48)
F_GREEN=RGBColor(0x4F,0x80,0x61); F_PURPLE=RGBColor(0x72,0x5C,0x91)
FONT="Pretendard"; FB="맑은 고딕"; EN="Arial"

L = json.load(io.open(sys.argv[1], encoding="utf-8"))
OUT = sys.argv[2]
P=L["passage"]; Q=L.get("question",{}); PTS=L.get("points",[]); CIRC="①②③④⑤"; TOTAL=7
ANS=int(L.get("answer") or Q.get("answer") or 1)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def setrun(r,t,sz,c,bold=False,ul=False,name=FONT,spc=None,strike=False):
    r.text=t; f=r.font; f.size=Pt(sz); f.bold=bold; f.underline=ul; f.color.rgb=c; f.name=name
    rPr=r._r.get_or_add_rPr(); ea=rPr.makeelement(qn('a:ea'),{'typeface':FB}); rPr.append(ea)
    if strike: rPr.set('strike','sngStrike')
    if spc is not None: rPr.set('spc',str(int(spc*100)))
def para(tf,t,sz,c,bold=False,align=PP_ALIGN.LEFT,after=4,before=0,name=FONT,first=False,spc=None,line=None):
    p=tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(after); p.space_before=Pt(before)
    if line: p.line_spacing=line
    setrun(p.add_run(),t,sz,c,bold,name=name,spc=spc); return p
def box(s,x,y,w,h):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tb.text_frame.word_wrap=True
    return tb,tb.text_frame
def rect(s,x,y,w,h,fill=None,ln=None,lw=1.0,round_=True,anchor=MSO_ANCHOR.TOP,nm=None):
    c=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: c.fill.background()
    else: c.fill.solid(); c.fill.fore_color.rgb=fill
    if ln is None: c.line.fill.background()
    else: c.line.color.rgb=ln; c.line.width=Pt(lw)
    c.shadow.inherit=False
    tf=c.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.16); tf.margin_right=Inches(0.14); tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.08)
    if nm: c.name=nm
    return c,tf
def newslide(bg):
    s=prs.slides.add_slide(BLANK)
    r,_=rect(s,0,0,13.333,7.5,bg,round_=False,nm="BG")
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s
def chrome(s,stage,title,page,dark=False,prompt=None,prompt_label=None):
    tcol=WHITE if dark else NAVY
    ch,tf=rect(s,0.45,0.30,1.95,0.36,(NAVY if not dark else GOLD),anchor=MSO_ANCHOR.MIDDLE)
    tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    setrun(p.add_run(),stage,11,(GOLDL if not dark else NAVY),True,name=EN,spc=1.4)
    _,tf=box(s,2.6,0.26,8.5,0.5); para(tf,title,21,tcol,True,first=True,after=0)
    _,tf=box(s,11.3,0.30,1.6,0.4)
    para(tf,page,12,GRAY if not dark else GOLDL,True,align=PP_ALIGN.RIGHT,first=True,after=0,name=EN)
    rect(s,0.45,0.78,12.44,0.016,GOLD if not dark else RGBColor(0x3a,0x4d,0x63),round_=False)
    if prompt:
        pb,tf=rect(s,0.45,6.52,12.44,0.44,(NAVY2 if dark else WHITE),(None if dark else BORDER),1.0,anchor=MSO_ANCHOR.MIDDLE)
        p=tf.paragraphs[0]
        setrun(p.add_run(),"  "+(prompt_label or "포인트")+"  ",10.5,GOLD,True,spc=1.2)
        setrun(p.add_run()," "+prompt,13,(WHITE if dark else INK))
    _,tf=box(s,0.45,7.08,6,0.3)
    para(tf,"이인혁 영어 · RAY ENGLISH",9,(RGBColor(0x6b,0x77,0x84) if dark else RGBColor(0xb9,0xb3,0xa4)),True,first=True,after=0,spc=1.4)
def clip(s,n): s=re.sub(r'\s+',' ',str(s or "")).strip(); return s if len(s)<=n else s[:n-1]+"…"
def run_html(p,text,sz,color=INK,circ_col=GOLD):
    ul=False
    for tok in re.split(r'(</?u>|[①②③④⑤])',str(text or "")):
        if tok=='<u>': ul=True; continue
        if tok=='</u>': ul=False; continue
        if not tok: continue
        if tok in CIRC: setrun(p.add_run(),tok,sz+1.5,circ_col,True,name=EN); continue
        clean=re.sub(r'<[^>]+>','',tok)
        if clean: setrun(p.add_run(),clean,sz,(F_RED if ul else color),bold=ul,ul=ul,name=EN)

# ── 1 표지 ──
s=newslide(NAVYD)
_,tf=box(s,0.9,0.5,11.6,0.4)
para(tf,f"{L['class']['school'].upper()} HIGH {L['class']['grade']} · {L['class']['lesson']}",13,GOLDL,True,first=True,name=EN,spc=1.6)
_,tf=box(s,0.9,1.75,11.6,2.0)
para(tf,P.get("title") or "어법 정확성 파악",34,WHITE,True,first=True,after=6,name=EN)
para(tf,clip(P.get("topic"),44),18,GOLDL,after=0)
c,tf=rect(s,0.9,4.15,11.5,1.35,NAVY2)
para(tf,"한 줄 요약",11,GOLD,True,first=True,after=3,spc=1.6)
para(tf,clip(P.get("one_liner"),74),15.5,WHITE,after=0,line=1.2)
_,tf=box(s,0.9,6.5,11.6,0.7)
para(tf,clip(L.get("source"),60),12,RGBColor(0x8a,0x97,0xa5),True,first=True,spc=0.8,after=2)
para(tf,f"이인혁 영어 · RAY ENGLISH · {L['class']['date']}",11,RGBColor(0x8a,0x97,0xa5),True,spc=1.2)

# ── 2 지문 통독 ──
s=newslide(IVORY); chrome(s,"READ","지문 통독 — 무슨 내용인가?",f"02 / {TOTAL}",prompt="어법을 보기 전에 글의 내용부터 잡는다")
c,tf=rect(s,0.45,1.0,12.44,4.75,WHITE,BORDER,1.2)
sz = 16 if len(P["text"])<=6 else 14.5
for i,t in enumerate(P["text"]):
    p=para(tf,"",1,INK,first=(i==0),after=7)
    setrun(p.add_run(),f"({i+1}) ",sz-3,GOLD,True,name=EN)
    setrun(p.add_run(),t,sz,INK,name=EN); p.line_spacing=1.22
c,tf=rect(s,0.45,5.9,12.44,0.5,IVORY,GOLD,1.2,anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]; setrun(p.add_run(),"요지  ",12,GOLD,True); setrun(p.add_run(),clip(P.get("one_liner"),96),13,NAVY,True)

# ── 3 실전 문항 ──
s=newslide(IVORY); chrome(s,"QUESTION","실전 문항 — 어법상 틀린 것은?",f"03 / {TOTAL}",prompt="밑줄마다 '무슨 문법인지' 먼저 명명하라")
c,tf=rect(s,0.45,1.0,1.35,0.42,NAVY,anchor=MSO_ANCHOR.MIDDLE)
tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p.add_run(),"어법",13,GOLDL,True,spc=1.2)
_,tf=box(s,1.95,0.97,10.9,0.5)
para(tf,Q.get("stem","다음 글의 밑줄 친 부분 중, 어법상 틀린 것은?"),16,INK,True,first=True)
c,tf=rect(s,0.45,1.6,9.8,4.75,WHITE,BORDER,1.1)
mp=Q.get("marked_passage","")
psz = 13 if len(re.sub(r'<[^>]+>','',mp))<820 else 11.5
p=para(tf,"",1,INK,first=True,after=0); p.line_spacing=1.32
run_html(p,mp,psz)
c,tf=rect(s,10.45,1.6,2.44,4.75,None,GRAY,1.0)
para(tf,"검산 판서",11,GRAY,True,first=True,align=PP_ALIGN.CENTER,after=6)
for lab in ["①","②","③","④","⑤"]:
    para(tf,lab+"  O / X",12.5,GRAY,after=9,name=EN)

# ── 4 밑줄 5포인트 진단 ──
s=newslide(IVORY); chrome(s,"DIAGNOSE","밑줄 5포인트 — 문법 항목 판정",f"04 / {TOTAL}",prompt="항목을 명명하면 정오는 저절로 드러난다")
y=1.05; H=min(1.0, 5.2/max(len(PTS),1))
for pt in PTS[:5]:
    isX = str(pt.get("verdict","")).upper()=="X"
    cl = F_RED if isX else F_GREEN
    # 번호 뱃지
    c,tf=rect(s,0.45,y,0.7,H,cl,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p.add_run(),CIRC[pt.get("n",1)-1],17,WHITE,True,name=EN)
    # 밑줄 표현
    c,tf=rect(s,1.25,y,2.5,H,WHITE,BORDER,1.0,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]
    setrun(p.add_run(),str(pt.get("underline","")),14.5,(F_RED if isX else INK),True,ul=isX,name=EN)
    if isX and pt.get("correct"):
        setrun(p.add_run(),"  → "+str(pt.get("correct")),13,F_GREEN,True,name=EN)
    # 문법 라벨
    c,tf=rect(s,3.85,y,2.9,H,IVORY,cl,1.2,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; setrun(p.add_run(),clip(pt.get("label"),22),13,cl,True)
    # 근거
    c,tf=rect(s,6.85,y,5.05,H,WHITE,BORDER,1.0,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; setrun(p.add_run(),clip(pt.get("why"),74),11.5,GRAY)
    # O/X
    c,tf=rect(s,12.0,y,0.89,H,cl,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p.add_run(),"X" if isX else "O",18,WHITE,True,name=EN)
    y+=H+0.06

# ── 5 정답·근거 ──
s=newslide(IVORY); chrome(s,"ANSWER","정답 · 결정적 근거",f"05 / {TOTAL}",prompt="틀린 형태 → 올바른 형태, 그리고 그 이유")
anspt=next((p for p in PTS if p.get("answer")), None)
c,tf=rect(s,0.45,1.05,2.3,1.7,None,GOLD,2.2,anchor=MSO_ANCHOR.MIDDLE)
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
setrun(p.add_run(),"정답\n",16,GRAY,True); setrun(p.add_run(),CIRC[ANS-1],40,NAVY,True,name=EN)
# 교정 카드
c,tf=rect(s,3.05,1.05,9.84,1.7,WHITE,BORDER,1.1,anchor=MSO_ANCHOR.MIDDLE)
if anspt:
    p=tf.paragraphs[0]
    setrun(p.add_run(),str(anspt.get("underline","")),22,F_RED,True,ul=True,strike=True,name=EN)
    setrun(p.add_run(),"   →   ",18,GRAY,True,name=EN)
    setrun(p.add_run(),str(anspt.get("correct","")),22,F_GREEN,True,name=EN)
    para(tf,clip(anspt.get("label"),40),13,GOLD,True,after=0,before=2)
c,tf=rect(s,0.45,3.0,12.44,1.5,WHITE,BORDER,1.1)
para(tf,"결정적 근거",11.5,GOLD,True,first=True,after=5,spc=1.2)
para(tf,clip(L.get("explanation_core"),210),13.5,INK,line=1.28)
# 나머지 정문 요약
c,tf=rect(s,0.45,4.7,12.44,1.6,IVORY,BORDER,1.1)
para(tf,"나머지는 왜 정문인가",11.5,GRAY,True,first=True,after=6,spc=1.0)
for pt in PTS:
    if pt.get("answer"): continue
    p=tf.add_paragraph(); p.space_after=Pt(4); p.line_spacing=1.15
    setrun(p.add_run(),f"{CIRC[pt.get('n',1)-1]} {pt.get('underline','')} ",12.5,F_GREEN,True,name=EN)
    setrun(p.add_run(),f"[{pt.get('label','')}] ",11.5,GOLD,True)
    setrun(p.add_run(),clip(pt.get("why"),64),11.5,INK)

# ── 6 직독직해 ──
s=newslide(IVORY); chrome(s,"READ","직독직해 — 문장별 해석",f"06 / {TOTAL}")
c,tf=rect(s,0.45,1.0,12.44,5.78,WHITE,BORDER,1.2)
ko=P.get("ko",[]); n=len(P["text"])
esz = 12.5 if n<=6 else (11 if n<=7 else 10)
aft = 7 if n<=6 else (4 if n<=7 else 3)
for i,t in enumerate(P["text"]):
    p=para(tf,"",1,INK,first=(i==0),after=1)
    setrun(p.add_run(),f"({i+1}) ",esz-2,GOLD,True,name=EN)
    setrun(p.add_run(),clip(t,240),esz,INK,name=EN); p.line_spacing=1.05
    if i<len(ko):
        p2=tf.add_paragraph(); p2.space_after=Pt(aft); p2.line_spacing=1.05
        setrun(p2.add_run(),"     "+clip(ko[i],150),esz-0.5,F_BLUE)

# ── 7 핵심 문법 · 어휘 · 과제 ──
s=newslide(NAVY); chrome(s,"CHECK","핵심 문법 · 어휘 · 과제",f"07 / {TOTAL}",dark=True,prompt_label="Exit",prompt="틀린 밑줄의 문법 항목을 한 문장으로 설명해 보라")
# 오답코드
wcs=(L.get("wrong_codes") or [])[:3]
x=0.45
for wc in wcs:
    c,tf=rect(s,x,1.08,4.05,1.5,NAVY2)
    para(tf,clip(wc.get("code"),14),16,GOLD,True,first=True,after=2)
    para(tf,clip(wc.get("miss"),24),12.5,WHITE,True,after=3)
    para(tf,"→ "+clip(wc.get("fix"),26),11.5,RGBColor(0xB9,0xC4,0xD2),after=0)
    x+=4.2
# 핵심 문법 3
c,tf=rect(s,0.45,2.85,7.7,3.45,None,GOLD,1.4)
para(tf,"핵심 문법 3",12,GOLD,True,first=True,after=7,spc=1.6)
for i,tx in enumerate((L.get("core3") or [])[:3]):
    p=tf.add_paragraph(); p.space_after=Pt(9); p.line_spacing=1.16
    setrun(p.add_run(),f"{i+1:02d} ",14,GOLD,True,name=EN); setrun(p.add_run(),clip(tx,58),13.5,WHITE)
# 어휘
c,tf=rect(s,8.35,2.85,4.54,3.45,NAVY2)
para(tf,"핵심 어휘",12,GOLDL,True,first=True,after=6,spc=1.6)
for v in (L.get("vocab") or [])[:8]:
    p=tf.add_paragraph(); p.space_after=Pt(5); p.line_spacing=1.1
    setrun(p.add_run(),clip(v.get("w"),20)+"  ",12,WHITE,True,name=EN)
    setrun(p.add_run(),clip(v.get("m"),20),11.5,RGBColor(0xC7,0xD0,0xDA))

prs.save(OUT)
print("SAVED", OUT, "/", len(prs.slides._sldIdLst), "slides")
