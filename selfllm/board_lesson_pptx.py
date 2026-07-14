# -*- coding: utf-8 -*-
"""board_lesson_pptx.py — 수업 JSON → RAY 보드 수업 PPTX 자동생성(기획서 레이아웃)
사용: python board_lesson_pptx.py <lesson.json> <out.pptx>"""
import sys, io, json, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY=RGBColor(0x1A,0x2A,0x3A); NAVY2=RGBColor(0x23,0x36,0x48); NAVYD=RGBColor(0x0F,0x1B,0x27)
GOLD=RGBColor(0xC5,0xA8,0x69); GOLDL=RGBColor(0xE8,0xD5,0xA5); IVORY=RGBColor(0xF7,0xF5,0xEF)
INK=RGBColor(0x23,0x2A,0x31); GRAY=RGBColor(0x5B,0x65,0x70); BORDER=RGBColor(0xE7,0xE3,0xD8); WHITE=RGBColor(0xFF,0xFF,0xFF)
F_BLUE=RGBColor(0x3E,0x7C,0xB1); F_ORANGE=RGBColor(0xD2,0x7A,0x3D); F_RED=RGBColor(0xB9,0x4A,0x48)
F_GREEN=RGBColor(0x4F,0x80,0x61); F_PURPLE=RGBColor(0x72,0x5C,0x91)
FONT="Pretendard"; FB="맑은 고딕"; EN="Arial"

L = json.load(io.open(sys.argv[1], encoding="utf-8"))
OUT = sys.argv[2]
P = L["passage"]; A = L["analysis"]; Q = L.get("question"); W = L.get("writing")
CIRC="①②③④⑤"; TOTAL=11

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def setrun(r,t,sz,c,bold=False,ul=False,name=FONT,spc=None):
    r.text=t; f=r.font; f.size=Pt(sz); f.bold=bold; f.underline=ul; f.color.rgb=c; f.name=name
    rPr=r._r.get_or_add_rPr(); ea=rPr.makeelement(qn('a:ea'),{'typeface':FB}); rPr.append(ea)
    if spc is not None: rPr.set('spc',str(int(spc*100)))
def para(tf,t,sz,c,bold=False,align=PP_ALIGN.LEFT,after=4,before=0,name=FONT,first=False,spc=None,line=None):
    p=tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment=align; p.space_after=Pt(after); p.space_before=Pt(before)
    if line: p.line_spacing=line
    setrun(p.add_run(),t,sz,c,bold,name=name,spc=spc); return p
def box(s,x,y,w,h):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tb.text_frame.word_wrap=True
    return tb,tb.text_frame
def rect(s,x,y,w,h,fill=None,ln=None,lw=1.0,round_=True,anchor=MSO_ANCHOR.TOP,nm=None):
    c=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: c.fill.background()
    else: c.fill.solid(); c.fill.fore_color.rgb=fill
    if ln is None: c.line.fill.background()
    else: c.line.color.rgb=ln; c.line.width=Pt(lw)
    c.shadow.inherit=False
    tf=c.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.16); tf.margin_right=Inches(0.14); tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.08)
    if nm: c.name=nm
    return c,tf
def newslide(bg):
    s=prs.slides.add_slide(BLANK)
    r,_=rect(s,0,0,13.333,7.5,bg,round_=False,nm="BG")
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s
def chrome(s,stage,title,page,dark=False,prompt_label=None,prompt=None):
    tcol=WHITE if dark else NAVY
    ch,tf=rect(s,0.45,0.30,1.85,0.36,(NAVY if not dark else GOLD),anchor=MSO_ANCHOR.MIDDLE,nm="HEADER_LABEL")
    tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    setrun(p.add_run(),stage,11,(GOLDL if not dark else NAVY),True,name=EN,spc=1.4)
    _,tf=box(s,2.5,0.26,8.6,0.5); para(tf,title,22,tcol,True,first=True,after=0)
    _,tf=box(s,11.3,0.30,1.6,0.4)
    para(tf,page,12,GRAY if not dark else GOLDL,True,align=PP_ALIGN.RIGHT,first=True,after=0,name=EN)
    rect(s,0.45,0.78,12.44,0.016,GOLD if not dark else RGBColor(0x3a,0x4d,0x63),round_=False)
    if prompt:
        pb,tf=rect(s,0.45,6.52,12.44,0.44,(NAVY2 if dark else WHITE),(None if dark else BORDER),1.0,anchor=MSO_ANCHOR.MIDDLE,nm="FOOTER_PROMPT")
        p=tf.paragraphs[0]
        setrun(p.add_run(),"  "+(prompt_label or "오늘의 질문")+"  ",10.5,GOLD,True,spc=1.2)
        setrun(p.add_run()," "+prompt,13,(WHITE if dark else INK))
    _,tf=box(s,0.45,7.08,4,0.3)
    para(tf,"이인혁 영어 · RAY ENGLISH",9,(RGBColor(0x6b,0x77,0x84) if dark else RGBColor(0xb9,0xb3,0xa4)),True,first=True,after=0,spc=1.4)

def wm(s):
    s=str(s or "")
    s=re.sub(r'\[Flow\$?\s*Edu\]\$?\s*flowedu\.tistory\.com','',s)
    s=re.sub(r'2027\$?.*?제작용\)','',s)
    return re.sub(r'\s+',' ',s).strip()
def clip(s,n): s=wm(s); return s if len(s)<=n else s[:n-1]+"…"

# ── 1 표지 ──
s=newslide(NAVYD)
_,tf=box(s,0.9,0.5,11,0.4)
para(tf,f"{L['class']['school'].upper()} HIGH {L['class']['grade']} · {L['class']['lesson']}",13,GOLDL,True,first=True,name=EN,spc=1.6)
_,tf=box(s,0.9,2.0,11.6,1.9)
para(tf,P.get("title") or "오늘의 지문",38,WHITE,True,first=True,after=6)
para(tf,clip(P.get("topic"),40),18,GOLDL,after=0)
c,tf=rect(s,0.9,4.4,11.5,1.15,NAVY2)
para(tf,"오늘의 질문",11,GOLD,True,first=True,after=3,spc=1.6)
para(tf,clip(A.get("conclusion") or A.get("main_idea_ko"),60),17,WHITE,after=0)
_,tf=box(s,0.9,6.6,11,0.35)
para(tf,f"이인혁 영어 · RAY ENGLISH · {L['class']['date']}",11,RGBColor(0x8a,0x97,0xa5),True,first=True,spc=1.2)

# ── 2 전체 지문 ──
s=newslide(IVORY); chrome(s,"READ","지문 전체 읽기",f"02 / {TOTAL}",prompt="이 글은 무엇에 대한 글인가?")
c,tf=rect(s,0.45,1.0,9.9,5.35,WHITE,BORDER,1.2,nm="PASSAGE_FULL")
sz = 21 if len(P["text"])<=6 else (18 if len(P["text"])<=8 else 16)
for i,t in enumerate(P["text"]):
    p=para(tf,"",1,INK,first=(i==0),after=8)
    setrun(p.add_run(),f"({i+1}) ",max(sz-6,12),GOLD,True,name=EN)
    setrun(p.add_run(),wm(t),sz,INK,name=EN); p.line_spacing=1.25
c,tf=rect(s,10.5,1.0,2.39,5.35,None,GRAY,1.0,nm="BOARD_AREA")
para(tf,"판서 여백",10.5,GRAY,True,first=True,align=PP_ALIGN.CENTER)

# ── 3 첫 대면 메모 ──
s=newslide(IVORY); chrome(s,"READ","첫 대면 메모 — 30초 판단",f"03 / {TOTAL}",prompt="전환 신호는 어느 단어인가?")
quads=[("소재","무엇에 관한 글인가?",A.get("material",""),F_BLUE),
       ("방향","긍정·부정·문제제기?",A.get("direction",""),F_RED),
       ("전환","흐름이 바뀌는 지점은?",A.get("shift",""),F_ORANGE),
       ("예상 결론","필자가 말하려는 것은?",clip(A.get("conclusion"),42),F_GREEN)]
pos=[(0.45,1.05),(6.75,1.05),(0.45,3.7),(6.75,3.7)]
for (lb,qq,ans,cl),(x,y) in zip(quads,pos):
    c,tf=rect(s,x,y,6.1,2.45,WHITE,cl,1.6)
    p=tf.paragraphs[0]; setrun(p.add_run(),lb,15,cl,True); p.space_after=Pt(2)
    para(tf,qq,12.5,GRAY,after=10)
    c2,tf2=rect(s,x+0.25,y+1.35,5.6,0.85,IVORY,GOLD,1.1,anchor=MSO_ANCHOR.MIDDLE,nm="ANSWER_MEMO")
    p=tf2.paragraphs[0]; setrun(p.add_run(),str(ans),14.5,NAVY,True)

# ── 4 핵심 문장 확대 ──
kn = int(P.get("key_sentence") or len(P["text"])); kn = min(max(kn,1),len(P["text"]))
key_en = wm(P["text"][kn-1])
lines = P.get("key_lines") or []
if not lines:
    parts = re.split(r'(?<=,)\s+', key_en)
    lines = parts if 2<=len(parts)<=4 else [key_en]
s=newslide(IVORY); chrome(s,"READ",f"핵심 문장 확대 — ({kn})",f"04 / {TOTAL}",prompt="주어는 어디서 끝나는가?")
_,tf=box(s,0.45,1.0,6,0.4)
para(tf,f"Sentence {kn}",13,GRAY,True,first=True,name=EN)
c,tf=rect(s,0.45,1.55,12.44,3.4,WHITE,BORDER,1.2,anchor=MSO_ANCHOR.MIDDLE,nm="SENTENCE_01")
ksz = 23 if len(key_en)<130 else 19
for i,ln in enumerate(lines):
    p=para(tf,ln,ksz,INK,first=(i==0),align=PP_ALIGN.CENTER,after=8,name=EN); p.line_spacing=1.3
_,tf=box(s,0.45,5.25,6.4,0.5)
para(tf,"Q. 이 문장의 동사는 무엇인가?",14.5,NAVY,True,first=True)
c,tf=rect(s,7.2,5.15,5.7,1.15,None,GRAY,1.0,nm="BOARD_AREA")
para(tf,"필기 구역",10.5,GRAY,True,first=True,align=PP_ALIGN.CENTER)

# ── 5 구조 해체 ──
comp = next((c for c in A.get("components",[]) if c.get("n")==kn), None)
def toplevel_segments(parse):
    """중첩 성분([M while [S humans] ...]) 안전 파싱: 최상위 라벨만 취하고 내부 라벨은 텍스트로 평탄화."""
    segs=[]; depth=0; start=0
    for i,ch in enumerate(parse):
        if ch=='[':
            if depth==0: start=i
            depth+=1
        elif ch==']':
            depth-=1
            if depth==0: segs.append(parse[start+1:i])
    out=[]
    for seg in segs:
        m=re.match(r'([SVOCM])\s+(.*)$',seg,re.S)
        if m:
            text=re.sub(r'\[([SVOCM])\s+','',m.group(2)).replace(']','').strip()
            out.append((m.group(1),re.sub(r'\s+',' ',text)))
    return out
segs = toplevel_segments(comp["parse"]) if comp else []
CMAP={"S":(F_BLUE,"S — 주어"),"V":(F_ORANGE,"V — 동사"),"O":(GOLD,"O — 목적어"),"C":(GOLD,"C — 보어"),"M":(F_PURPLE,"M — 수식·부가")}
s=newslide(IVORY); chrome(s,"STRUCTURE","구조 해체 — 성분으로 끊기",f"05 / {TOTAL}",prompt="뼈대(S·V·O)만 남기면 무엇이 남는가?")
totlen=sum(len(t) for _,t in segs) or len(key_en)
sfz = 15.5 if totlen<110 else (13.5 if totlen<160 else 12)
c,tf=rect(s,0.45,1.05,12.44,1.2,WHITE,BORDER,1.2,anchor=MSO_ANCHOR.MIDDLE,nm="SENTENCE_01")
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
if segs:
    for lab,txt in segs:
        cl=CMAP.get(lab,(INK,""))[0]
        setrun(p.add_run(),txt+" ",sfz,cl,True,ul=(lab in "SV"),name=EN)
else:
    setrun(p.add_run(),key_en,sfz,INK,name=EN)
n=min(len(segs),5) or 1
rh=min(0.6, 3.1/n); y=2.5
for lab,txt in segs[:5]:
    cl,lbl=CMAP.get(lab,(GRAY,lab))
    c,tf=rect(s,0.45,y,3.3,rh,cl,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; setrun(p.add_run(),lbl,12,WHITE,True)
    c,tf=rect(s,3.9,y,8.99,rh,WHITE,BORDER,1.0,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; setrun(p.add_run(),clip(txt,80),12.5,INK,name=EN)
    y+=rh+0.12
c,tf=rect(s,0.45,5.82,12.44,0.6,IVORY,GOLD,1.2,anchor=MSO_ANCHOR.MIDDLE,nm="ANSWER_TRANSLATION")
p=tf.paragraphs[0]
setrun(p.add_run(),"해석  ",12,GOLD,True)
setrun(p.add_run(),clip(P.get("ko",[""]*kn)[kn-1] if len(P.get("ko",[]))>=kn else "",90),13.5,NAVY,True)

# ── 6 결속 ──
s=newslide(IVORY); chrome(s,"FLOW","결속 — 같은 단어 없이 다시 가리키기",f"06 / {TOTAL}",prompt="이 표현은 앞의 무엇을 다시 가리키는가?")
coh = A.get("cohesion") or []
c,tf=rect(s,0.45,1.05,5.6,4.1,WHITE,BORDER,1.2)
para(tf,"앞 표현",12,GRAY,True,first=True,name=EN,after=8)
for pair in coh[:3]:
    p=tf.add_paragraph(); p.space_after=Pt(14); p.line_spacing=1.25
    setrun(p.add_run(),clip(pair.get("from",""),44),16,F_PURPLE,True,ul=True,name=EN)
c,tf=rect(s,7.3,1.05,5.6,4.1,WHITE,BORDER,1.2)
para(tf,"재지시 표현",12,GRAY,True,first=True,name=EN,after=8)
for pair in coh[:3]:
    p=tf.add_paragraph(); p.space_after=Pt(14); p.line_spacing=1.25
    setrun(p.add_run(),clip(pair.get("to",""),44),16,F_PURPLE,True,ul=True,name=EN)
for i in range(min(len(coh),3)):
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(6.12),Inches(1.95+i*0.95),Inches(1.1),Inches(0.16))
    a.fill.solid(); a.fill.fore_color.rgb=F_PURPLE; a.line.fill.background(); a.shadow.inherit=False
if coh:
    chain=" → ".join([coh[0].get("from","")]+[c2.get("to","") for c2 in coh[:3]])
    c,tf=rect(s,0.45,5.4,12.44,0.9,None,F_PURPLE,1.3,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; setrun(p.add_run(),clip(chain,110),13.5,F_PURPLE,True,name=EN)

# ── 7 흐름 지도 ──
s=newslide(IVORY); chrome(s,"FLOW","단락 흐름 지도",f"07 / {TOTAL}",prompt="어느 문장이 결론인가?")
def rolecolor(r):
    r=str(r)
    if re.search(r"도입|통념|배경",r): return NAVY2
    if re.search(r"전환|대조|반박",r): return F_RED
    if re.search(r"주장|핵심|논지",r): return F_BLUE
    if re.search(r"근거|예시|부연|설명",r): return F_GREEN
    if re.search(r"결론|함의|재진술",r): return GOLD
    return GRAY
flows=(A.get("logic") or [])[:5]
y=1.1; H=min(0.92, 4.4/max(len(flows),1))
for i,f in enumerate(flows):
    cl=rolecolor(f.get("role"))
    c,tf=rect(s,2.2,y,8.9,H,None,cl,1.7,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]
    setrun(p.add_run(),f" ({f.get('n')}) {f.get('role','')}  ",13.5,cl,True)
    setrun(p.add_run(),clip(f.get("note","") or (P["ko"][f.get("n",1)-1] if len(P.get("ko",[]))>=f.get("n",1) else ""),52),13,INK)
    if i<len(flows)-1:
        a=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(6.5),Inches(y+H+0.02),Inches(0.25),Inches(0.14))
        a.fill.solid(); a.fill.fore_color.rgb=GOLD; a.line.fill.background(); a.shadow.inherit=False
    y+=H+0.2
c,tf=rect(s,0.45,5.85,12.44,0.55,IVORY,GOLD,1.1,anchor=MSO_ANCHOR.MIDDLE,nm="ANSWER_CORE")
p=tf.paragraphs[0]
setrun(p.add_run(),"주장  ",12,GOLD,True); setrun(p.add_run(),clip(A.get("main_idea_ko"),80),14,NAVY,True)

# ── 8 실전 문제 ──
s=newslide(IVORY); chrome(s,"QUESTION","실전 문제",f"08 / {TOTAL}",prompt="대상·범위·방향·관계·강도를 검산하라")
def run_html(p,text,sz,color=INK):
    ul=False
    for tok in re.split(r'(</?u>)',str(text or "")):
        if tok=='<u>': ul=True; continue
        if tok=='</u>': ul=False; continue
        if not tok: continue
        clean=re.sub(r'<[^>]+>','',tok)
        if clean: setrun(p.add_run(),clean,sz,(F_RED if ul else color),ul,name=EN)
if Q:
    c,tf=rect(s,0.45,1.0,1.35,0.4,NAVY,anchor=MSO_ANCHOR.MIDDLE)
    tf.word_wrap=False; p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; setrun(p.add_run(),"CORE",11,GOLDL,True,name=EN,spc=1.2)
    _,tf=box(s,1.95,0.97,10.9,0.5)
    para(tf,Q.get("stem",""),16.5,INK,True,first=True)
    c,tf=rect(s,0.45,1.62,9.7,4.75,WHITE,BORDER,1.1,nm="PASSAGE_FULL")
    ch=Q.get("choices",[])[:5]
    labelish=ch and all(len(str(x).strip())<=3 for x in ch)
    qpg=wm(Q.get("passage") or "")
    if qpg:
        psz = 12.5 if len(qpg)<900 else 11
        p=para(tf,"",1,INK,first=True,after=10); p.line_spacing=1.3
        run_html(p,qpg,psz)
    if labelish:
        p=tf.add_paragraph(); p.space_before=Pt(6)
        setrun(p.add_run(),"   ".join(CIRC[i]+" "+str(o) for i,o in enumerate(ch)),14.5,INK,True,name=EN)
    else:
        for i,o in enumerate(ch):
            p=para(tf,"",1,INK,first=(not qpg and i==0),after=10)
            setrun(p.add_run(),CIRC[i]+"  ",14.5,INK,True,name=EN)
            setrun(p.add_run(),str(o),14.5,INK,name=EN); p.line_spacing=1.2
    c,tf=rect(s,10.35,1.62,2.55,4.75,None,GRAY,1.0,nm="BOARD_AREA")
    para(tf,"검산 판서",11,GRAY,True,first=True,align=PP_ALIGN.CENTER)

# ── 9 정답·근거 ──
s=newslide(IVORY); chrome(s,"ANSWER","정답·근거",f"09 / {TOTAL}",prompt="오답의 '어긋난 지점'만 짚는다")
if Q:
    exp=str(Q.get("explanation",""))
    core = exp.split("[선지 분석]")[0].split("【")[0].split("[오답")[0].strip()
    c,tf=rect(s,0.45,1.05,2.1,1.5,None,GOLD,2.2,anchor=MSO_ANCHOR.MIDDLE,nm="ANSWER_REVEAL")
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    setrun(p.add_run(),"정답\n"+CIRC[Q.get("answer",1)-1],26,NAVY,True)
    c,tf=rect(s,2.85,1.05,10.04,1.5,WHITE,BORDER,1.1,nm="ANSWER_EVIDENCE")
    para(tf,"결정적 근거",11.5,GOLD,True,first=True,after=4,spc=1.2)
    para(tf,clip(core,170),13.5,INK,line=1.25)
    c,tf=rect(s,0.45,2.85,12.44,3.45,WHITE,BORDER,1.1,nm="ANSWER_WHYNOT")
    para(tf,"오답 검산",12,GRAY,True,first=True,after=8)
    m=re.search(r"\[선지 분석\]([^\n【]*)",exp)
    items=[]
    if m: items=[x.strip() for x in re.split(r"·|·",m.group(1)) if x.strip()][:5]
    if not items:
        items=[x.strip() for x in re.findall(r"[①②③④⑤][^\n①②③④⑤]{6,120}",exp.split("[오답 진단·처방]")[-1])][:4]
    for it in items:
        p=tf.add_paragraph(); p.space_after=Pt(7); setrun(p.add_run(),clip(it,110),13,GRAY)

# ── 10 조건부 서술형 ──
s=newslide(IVORY); chrome(s,"WRITE","조건부 서술형",f"10 / {TOTAL}",prompt="배열 전에 조건부터 체크")
if W:
    ins=str(W.get("instruction",""))
    exp=str(W.get("explanation",""))
    model=""
    if "[모범답안]" in exp: model=exp.split("[모범답안]")[1].split("(채점")[0].split("[채점")[0].split("【")[0].strip()
    c,tf=rect(s,0.45,1.0,4.9,5.35,NAVY)
    para(tf,"문제·조건",12,GOLDL,True,first=True,after=8,spc=1.6)
    for lnn in [x for x in ins.split("\n") if x.strip()][:8]:
        para(tf,clip(lnn,46),12.5,WHITE,after=6,line=1.2)
    c,tf=rect(s,5.6,1.0,7.29,4.15,WHITE,BORDER,1.2,nm="BOARD_AREA")
    para(tf,"문장 제작 보드",12,GRAY,True,first=True,after=8)
    for st,cl in [("① 의미 뼈대",INK),("② 주어 선택",F_BLUE),("③ 동사 형태",F_ORANGE),("④ 배열",F_RED),("⑤ 검산(단어수·수일치)",F_GREEN)]:
        p=tf.add_paragraph(); p.space_after=Pt(13); setrun(p.add_run(),st,14.5,cl,True)
    c,tf=rect(s,5.6,5.35,7.29,1.0,IVORY,GOLD,1.4,anchor=MSO_ANCHOR.MIDDLE,nm="ANSWER_MODEL")
    p=tf.paragraphs[0]
    setrun(p.add_run(),"모범  ",12,GOLD,True); setrun(p.add_run(),clip(model,120),12.5,NAVY,True,name=EN)

# ── 11 오답코드·핵심3·과제 ──
s=newslide(NAVY); chrome(s,"CHECK","오답코드 · 오늘의 핵심 · 과제",f"11 / {TOTAL}",dark=True,prompt_label="Exit",prompt="오늘 주장을 한 문장으로 쓰고 자신의 오답코드를 고르시오")
wcs=(L.get("wrong_codes") or [])[:3]
if not wcs: wcs=[{"code":"F","miss":"흐름·결속","fix":"지시어 추적"},{"code":"S","miss":"긴 문장","fix":"절 경계 훈련"},{"code":"WR","miss":"조건 누락","fix":"체크 후 재작성"}]
x=0.45
for wc in wcs:
    c,tf=rect(s,x,1.1,4.05,1.6,NAVY2)
    para(tf,wc.get("code",""),20,GOLD,True,first=True,after=2,name=EN)
    para(tf,clip(wc.get("miss"),22),13,WHITE,True,after=3)
    para(tf,"→ "+clip(wc.get("fix"),24),11.5,RGBColor(0xB9,0xC4,0xD2),after=0)
    x+=4.2
c,tf=rect(s,0.45,3.0,12.44,1.8,None,GOLD,1.4)
para(tf,"오늘의 핵심 3",12,GOLD,True,first=True,after=6,spc=1.6)
for i,tx in enumerate((L.get("core3") or [])[:3]):
    p=tf.add_paragraph(); p.space_after=Pt(5)
    setrun(p.add_run(),f"{i+1:02d} ",14,GOLD,True,name=EN); setrun(p.add_run(),clip(tx,70),14,WHITE)
c,tf=rect(s,0.45,5.0,12.44,1.35,NAVY2)
para(tf,"과제",12,GOLDL,True,first=True,after=5,spc=1.6)
para(tf,"필수  1. 지문 3회 음독  2. 핵심 문장 구조 재표시  3. 변형 예상 2문항 만들기",13.5,WHITE,after=3)
para(tf,"오답코드별  "+" · ".join(f"{w.get('code')}: {clip(w.get('fix'),18)}" for w in wcs),13.5,GOLDL,after=0)

prs.save(OUT)
print("SAVED", OUT, "/", len(prs.slides._sldIdLst), "slides")
