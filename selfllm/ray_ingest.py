# -*- coding: utf-8 -*-
"""ray_ingest.py — 다양한 파일에서 지문(텍스트) 추출
지원: PDF · DOCX · PPTX · HWP/HWPX · XLSX · PNG/JPG(로컬 OCR) · TXT
사용:
  python ray_ingest.py <파일경로>            # 추출 텍스트 출력
  from ray_ingest import ingest; ingest(path) # 라이브러리
이미지와 스캔 PDF는 외부 전송 없이 프로젝트 내 오프라인 OCR로 인식한다."""
import sys, io, os, re, zipfile, hashlib
HERE = os.path.dirname(os.path.abspath(__file__))
CACHE_DIR = os.path.join(os.environ.get("LOCALAPPDATA", os.path.expanduser("~")), "RAYEnglish", "ingest_cache")

def _clean(t): return re.sub(r"[ \t]+\n", "\n", re.sub(r"\n{3,}", "\n\n", (t or "").strip()))

def read_txt(p):
    raw = io.open(p, "rb").read()
    for encoding in ("utf-8-sig", "utf-16", "cp949", "latin-1"):
        try: return raw.decode(encoding)
        except UnicodeDecodeError: pass
    return raw.decode("utf-8", "replace")

def read_docx(p):
    import docx
    d = docx.Document(p)
    parts = [para.text for para in d.paragraphs]
    for tbl in d.tables:
        for row in tbl.rows: parts.append("\t".join(c.text for c in row.cells))
    for section in d.sections:
        parts.extend(para.text for para in section.header.paragraphs if para.text)
        parts.extend(para.text for para in section.footer.paragraphs if para.text)
    return "\n".join(parts)

def read_pptx(p):
    from pptx import Presentation
    prs = Presentation(p); out = []
    def shape_text(sh):
        values = []
        if getattr(sh, "has_text_frame", False):
            values.extend("".join(r.text for r in para.runs) or para.text for para in sh.text_frame.paragraphs)
        if getattr(sh, "has_table", False):
            for row in sh.table.rows:
                values.append("\t".join(cell.text for cell in row.cells))
        if getattr(sh, "shape_type", None) == 6:
            for child in sh.shapes:
                values.extend(shape_text(child))
        return [value for value in values if value]
    for number, s in enumerate(prs.slides, 1):
        out.append(f"[Slide {number}]")
        for sh in s.shapes:
            out.extend(shape_text(sh))
        try:
            notes = s.notes_slide.notes_text_frame.text.strip()
            if notes:
                out.append("[Notes]\n" + notes)
        except Exception:
            pass
    return "\n".join(x for x in out if x)

def read_hwp(p):
    """HWP: 미리보기 텍스트(PrvText) 스트림 추출(olefile). 대부분의 지문 추출에 충분."""
    import olefile
    ole = olefile.OleFileIO(p)
    try:
        if ole.exists("PrvText"):
            return ole.openstream("PrvText").read().decode("utf-16-le", errors="ignore")
        # 폴백: 텍스트 유사 스트림 긁기
        chunks = []
        for entry in ole.listdir():
            name = "/".join(entry)
            if "Text" in name or "Body" in name:
                try: chunks.append(ole.openstream(entry).read().decode("utf-16-le", errors="ignore"))
                except Exception: pass
        return "\n".join(chunks)
    finally:
        ole.close()

def read_image(p):
    from ray_local_ocr import ocr_image
    return ocr_image(p)

def read_pdf(p):
    import fitz
    from ray_local_ocr import ocr_image

    max_ocr_pages = max(1, min(int(os.environ.get("RAY_OCR_MAX_PDF_PAGES", "12")), 40))
    out = []
    document = fitz.open(p)
    try:
        page_count = document.page_count
        if page_count <= max_ocr_pages:
            ocr_pages = set(range(page_count))
        elif max_ocr_pages == 1:
            ocr_pages = {0}
        else:
            ocr_pages = {
                round(index * (page_count - 1) / (max_ocr_pages - 1))
                for index in range(max_ocr_pages)
            }
        for index, page in enumerate(document):
            text = page.get_text("text").strip()
            if len(text) >= 30:
                out.append(f"[Page {index + 1}]\n{text}")
            elif index in ocr_pages:
                try:
                    pixmap = page.get_pixmap(dpi=180, alpha=False)
                    recognized = ocr_image(pixmap.tobytes("png"))
                    if recognized:
                        out.append(f"[Page {index + 1} · Local OCR]\n{recognized}")
                except Exception as exc:
                    out.append(f"[Page {index + 1} · OCR 실패] {exc}")
        if page_count > max_ocr_pages:
            out.append(f"[OCR 안내] 스캔 페이지는 전체 {page_count}쪽 중 {len(ocr_pages)}쪽을 고르게 표본 인식함")
    finally:
        document.close()
    return "\n\n".join(out)

def read_hwpx(p):
    parts = []
    with zipfile.ZipFile(p) as archive:
        names = [name for name in archive.namelist() if name.lower().endswith(".xml") and "section" in name.lower()]
        for name in sorted(names):
            raw = archive.read(name).decode("utf-8", "replace")
            parts.extend(re.findall(r">([^<>]+)<", raw))
    return "\n".join(parts)

def read_xlsx(p):
    from openpyxl import load_workbook
    workbook = load_workbook(p, read_only=True, data_only=True)
    out = []
    try:
        for sheet in workbook.worksheets:
            out.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value is not None]
                if values:
                    out.append("\t".join(values))
    finally:
        workbook.close()
    return "\n".join(out)

READERS = {".txt": read_txt, ".md": read_txt, ".csv": read_txt, ".json": read_txt,
           ".docx": read_docx, ".pptx": read_pptx, ".hwp": read_hwp,
           ".hwpx": read_hwpx, ".xlsx": read_xlsx, ".pdf": read_pdf,
           ".png": read_image, ".jpg": read_image, ".jpeg": read_image,
           ".bmp": read_image, ".webp": read_image, ".tif": read_image,
           ".tiff": read_image}

def ingest(path):
    path = os.path.abspath(path)
    ext = os.path.splitext(path)[1].lower()
    fn = READERS.get(ext)
    if not fn: raise ValueError(f"지원하지 않는 형식: {ext} (지원: {', '.join(READERS)})")
    stat = os.stat(path)
    key = hashlib.sha256(f"{path}|{stat.st_size}|{stat.st_mtime_ns}|v3".encode("utf-8")).hexdigest()
    cache_path = os.path.join(CACHE_DIR, key + ".txt")
    try:
        if os.path.isfile(cache_path):
            return io.open(cache_path, encoding="utf-8").read()
    except OSError:
        pass
    text = _clean(fn(path))
    try:
        os.makedirs(CACHE_DIR, exist_ok=True)
        io.open(cache_path, "w", encoding="utf-8").write(text)
    except OSError:
        pass
    return text

def extract_english_passage(text, min_words=25):
    """추출 텍스트에서 가장 그럴듯한 '영어 지문' 한 덩어리를 고름(문항용)."""
    # 영어 비중 높은 문단 우선
    paras = re.split(r"\n\s*\n", text)
    scored = []
    for para in paras:
        eng = len(re.findall(r"[A-Za-z]", para)); tot = len(re.sub(r"\s", "", para)) or 1
        words = len(re.findall(r"[A-Za-z]+", para))
        if words >= min_words and eng / tot > 0.6: scored.append((words, para.strip()))
    if scored:
        scored.sort(reverse=True); return re.sub(r"\s+", " ", scored[0][1])
    return re.sub(r"\s+", " ", text)[:2000]

if __name__ == "__main__":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
    if len(sys.argv) < 2: print(__doc__); sys.exit()
    t = ingest(sys.argv[1])
    print("=== 추출 텍스트 (앞 800자) ===\n" + t[:800])
    print("\n=== 자동 선별 영어 지문 ===\n" + extract_english_passage(t)[:800])
