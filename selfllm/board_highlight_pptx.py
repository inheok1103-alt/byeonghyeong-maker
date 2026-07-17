# -*- coding: utf-8 -*-
"""board_highlight_pptx.py — 판서형(검정 배경·색 하이라이트) 수업 PPTX 렌더러
사진 스타일 차용: 배경지식 hook → 문제 → 전면 하이라이트 분석 → 흐름 → 정답.
연결어/핵심개념/구조어를 색으로, 어려운 단어는 인라인 뜻풀이, 축약(…) 절대 없음.
사용: python board_highlight_pptx.py <data.json> <out.pptx>
kind: "reading" | "grammar"   마크업: <<TEXT|hl=color,tc=color,g=뜻,u,it,co,cc>>"""
import sys, io, json, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

def H(x): return RGBColor(int(x[0:2],16),int(x[2:4],16),int(x[4:6],16))
BG=H("0B0C10"); PANEL=H("15171E"); PANEL2=H("1D2029"); LINE=H("2A2E3A")
WHITE=H("F2F3F5"); MUTE=H("9AA3B0"); GOLD=H("D8B968"); GOLDL=H("E8D5A5"); DARK=H("14151A")
HL={"red":(H("D93B3B"),WHITE),"magenta":(H("D0479A"),WHITE),"pink":(H("E86BB0"),DARK),
    "blue":(H("2E7DD1"),WHITE),"cyan":(H("28C2C8"),DARK),"green":(H("39B968"),DARK),
    "yellow":(H("F0C22E"),DARK),"orange":(H("EE8A45"),DARK),"purple":(H("8E67D6"),WHITE)}
TC={"red":H("FF6B6B"),"magenta":H("F072BE"),"pink":H("F090C6"),"blue":H("6FB3F0"),"cyan":H("4FD3D8"),
    "green":H("6FE39A"),"yellow":H("F4CE5A"),"orange":H("F0A868"),"purple":H("C6A2F5"),
    "white":WHITE,"gold":GOLDL,"mute":MUTE}
SW={"green":H("39B968"),"red":H("D93B3B"),"magenta":H("D0479A"),"blue":H("2E7DD1"),"gold":GOLD,
    "yellow":H("F0C22E"),"orange":H("EE8A45"),"purple":H("8E67D6"),"cyan":H("28C2C8")}
EN="Segoe UI"; ENB="Segoe UI Semibold"; KO="맑은 고딕"; CIRC="①②③④⑤⑥⑦⑧⑨⑩"

L=json.load(io.open(sys.argv[1],encoding="utf-8")); OUT=sys.argv[2]
KIND=L.get("kind","reading")
prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def _ea(r,name=KO):
    rPr=r._r.get_or_add_rPr(); ea=rPr.makeelement(qn('a:ea'),{'typeface':name}); rPr.append(ea)
def _hl(r,rgb):
    rPr=r._r.get_or_add_rPr()
    h=rPr.makeelement(qn('a:highlight'),{}); c=rPr.makeelement(qn('a:srgbClr'),{'val':'%02X%02X%02X'%(rgb[0],rgb[1],rgb[2])}); h.append(c)
    latin=rPr.find(qn('a:latin')); ea=rPr.find(qn('a:ea')); ref=latin if latin is not None else ea
    (ref.addprevious(h) if ref is not None else rPr.append(h))
def run(p,t,sz,c,bold=False,ul=False,it=False,name=EN,hl=None,spc=None):
    r=p.add_run(); r.text=t; f=r.font; f.size=Pt(sz); f.bold=bold; f.underline=ul; f.italic=it
    f.color.rgb=c; f.name=name; rPr=r._r.get_or_add_rPr()
    if spc is not None: rPr.set('spc',str(int(spc*100)))
    _ea(r)
    if hl is not None: _hl(r,hl)
    return r
def box(s,x,y,w,h,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.12); tf.margin_right=Inches(0.12); tf.margin_top=Inches(0.05); tf.margin_bottom=Inches(0.05)
    return tf
def rect(s,x,y,w,h,fill=None,ln=None,lw=1.0,anchor=MSO_ANCHOR.TOP,round_=True):
    c=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: c.fill.background()
    else: c.fill.solid(); c.fill.fore_color.rgb=fill
    if ln is None: c.line.fill.background()
    else: c.line.color.rgb=ln; c.line.width=Pt(lw)
    c.shadow.inherit=False; tf=c.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.16); tf.margin_right=Inches(0.14); tf.margin_top=Inches(0.10); tf.margin_bottom=Inches(0.10)
    return c,tf
def slide():
    s=prs.slides.add_slide(BLANK); r,_=rect(s,0,0,13.333,7.5,BG,round_=False)
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element); return s
def head(s,tag,title,page):
    c,tf=rect(s,0.5,0.34,2.05,0.44,GOLD,anchor=MSO_ANCHOR.MIDDLE); tf.word_wrap=False
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run(p,tag,12.5,DARK,True,name=ENB,spc=1.2)
    tf=box(s,2.7,0.30,9.0,0.6,MSO_ANCHOR.MIDDLE); run(tf.paragraphs[0],title,23,WHITE,True,name=KO)
    tf=box(s,11.4,0.36,1.5,0.4,MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT; run(p,page,12,MUTE,True,name=ENB)
    rect(s,0.5,0.92,12.33,0.02,GOLD,round_=False)
def foot(s,note=None):
    tf=box(s,0.5,7.08,7,0.3); run(tf.paragraphs[0],"이인혁 영어 · RAY ENGLISH",9,H("5B6470"),True,name=ENB,spc=1.2)
    if note:
        tf=box(s,6.5,7.08,6.33,0.3,MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.RIGHT; run(p,note,10.5,MUTE,name=KO)

def parse_attrs(a):
    d={}
    for tok in a.split(","):
        tok=tok.strip()
        if not tok: continue
        if "=" in tok: k,v=tok.split("=",1); d[k.strip()]=v.strip()
        else: d[tok]=True
    return d
def flatten_markup(text):
    """중첩 마크업(<<A <<B|g=뜻>>|tc=..>>) 방어: 바깥 태그 '안'에 있는 안쪽 태그를 평문(+뜻풀이)으로 녹인다."""
    s=str(text or "")
    for _ in range(6):
        changed=[False]
        def repl(m):
            before=s[:m.start()]
            if before.count("<<") > before.count(">>"):      # 바깥 태그 안 = 중첩
                changed[0]=True
                a=parse_attrs(m.group(2)); g=a.get("g")
                return m.group(1)+(("("+str(g)+")") if g and g is not True else "")
            return m.group(0)
        s2=re.sub(r'<<([^<>]*?)\|([^<>]*?)>>', repl, s)
        s=s2
        if not changed[0]: break
    if s.count("<<") != s.count(">>"):                        # 짝 안 맞으면 태그 제거(누수 방지)
        s=s.replace("<<","").replace(">>","")
    return s

def emit(p,text,sz,base_c=WHITE,base_name=EN):
    text=flatten_markup(text)
    for tok in re.split(r'(<<.*?>>)',text):
        if not tok: continue
        if tok.startswith("<<") and tok.endswith(">>"):
            inner=tok[2:-2]; txt,_,att=inner.partition("|"); a=parse_attrs(att)
            c=base_c; hl=None; bold=False; ul="u" in a; it="it" in a
            if "co" in a: hl=HL["green"][0]; c=HL["green"][1]; bold=True
            elif "cc" in a: c=GOLDL; bold=True
            if "hl" in a and a["hl"] in HL: hl=HL[a["hl"]][0]; c=HL[a["hl"]][1]
            if "tc" in a and a["tc"] in TC: c=TC[a["tc"]]
            run(p,txt,sz,c,bold=bold,ul=ul,it=it,name=base_name,hl=hl)
            if "g" in a and a["g"] is not True:
                run(p,"("+str(a["g"])+")",sz*0.84,GOLDL,name=KO)
        else:
            run(p,tok,sz,base_c,name=base_name)
def strip_markup(text):
    text=flatten_markup(text)
    out=[]
    for tok in re.split(r'(<<.*?>>)',text):
        if tok.startswith("<<") and tok.endswith(">>"):
            inner=tok[2:-2]; txt,_,att=inner.partition("|"); a=parse_attrs(att); out.append(txt)
            if "g" in a and a["g"] is not True: out.append("("+str(a["g"])+")")
        else: out.append(tok)
    return "".join(out)
def plen(t): return len(strip_markup(t))

# 페이지 계산
NPASS = 2 if plen(L["passage"])>1500 else 1
TOTAL = (1 if L.get("hook") else 0) + 1 + NPASS + (2 if L.get("sents") and len(L["sents"])>8 else (1 if L.get("sents") else 0)) + (1 if L.get("flow") else 0) + 1
PG=[0]
def pg():
    PG[0]+=1; return f"{PG[0]:02d} / {TOTAL:02d}"

# ═══ HOOK : 배경지식으로 흥미 유발 ═══
if L.get("hook"):
    hk=L["hook"]; s=slide(); head(s,"HOOK",hk.get("title","오늘의 배경지식"),pg())
    tf=box(s,0.55,1.06,12.2,0.72); p=tf.paragraphs[0]; p.line_spacing=1.22; emit(p,hk.get("lead",""),15.5,GOLDL)
    cards=hk.get("cards",[])[:4]; cw=6.06; gap=0.22; x0=0.55; y0=1.92; ch=1.55
    for i,cd in enumerate(cards):
        x=x0+(i%2)*(cw+gap); y=y0+(i//2)*(ch+0.2); cl=SW.get(cd.get("cl","gold"),GOLD)
        cc,tf=rect(s,x,y,cw,ch,PANEL,cl,1.6)
        run(tf.paragraphs[0],cd.get("head",""),14.5,cl,True,name=KO)
        p=tf.add_paragraph(); p.space_before=Pt(3); p.line_spacing=1.24; emit(p,cd.get("body",""),12.5,H("D5DAE2"))
    rows=(len(cards)+1)//2; by=y0+rows*(ch+0.2)+0.02
    cc,tf=rect(s,0.55,min(by,5.72),12.28,0.92,GOLD,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; run(p,"오늘의 질문   ",12.5,DARK,True,name=KO,spc=1.0); run(p,hk.get("bridge",""),15,DARK,True,name=KO)
    foot(s,"배경지식으로 예열 → 지문으로 진입")

# ═══ 문제 ═══
s=slide(); head(s,"Q",L.get("header","실전 문제"),pg())
tf=box(s,0.55,1.02,9,0.36); run(tf.paragraphs[0],L.get("meta",""),12.5,MUTE,name=KO,spc=0.4)
stem=L.get("stem", "다음 글의 빈칸에 들어갈 말로 가장 적절한 것은?" if KIND=="reading" else "다음 글의 밑줄 친 부분 중, 어법상 틀린 것은?")
tf=box(s,0.55,1.36,12.2,0.5); run(tf.paragraphs[0],stem,16.5,GOLDL,True,name=KO)
if KIND=="reading":
    pc,ptf=rect(s,0.5,1.94,12.33,3.32,PANEL,LINE,1.0)
    _pl=plen(L["passage"])
    psz=15 if _pl<1150 else (13 if _pl<1450 else (11.5 if _pl<1750 else 10.5))
    p=ptf.paragraphs[0]; p.line_spacing=(1.3 if _pl<1450 else 1.18); emit(p,L["passage"],psz,WHITE)
    ctf=box(s,0.62,5.36,12.15,1.5)
    for i,c in enumerate(L.get("choices",[])):
        p=ctf.paragraphs[0] if i==0 else ctf.add_paragraph(); p.space_after=Pt(2); p.line_spacing=1.1
        run(p,CIRC[c.get("n",i+1)-1]+" ",13.5,GOLDL,True,name=EN); emit(p,c.get("text",""),13.5,WHITE)
else:
    pc,ptf=rect(s,0.5,1.94,12.33,4.25,PANEL,LINE,1.0)
    psz=15.5 if plen(L["passage"])<1150 else 14.5
    p=ptf.paragraphs[0]; p.line_spacing=1.34; emit(p,L["passage"],psz,WHITE)
    tf=box(s,0.55,6.32,12.2,0.4); run(tf.paragraphs[0],"밑줄 ①~⑤ 중 어법상 틀린 것 하나를 고르시오.",12.5,MUTE,it=True,name=KO)
foot(s,"먼저 스스로 풀고, 연결어에 표시하며 흐름을 잡는다")

# ═══ 지문 분석(전면 하이라이트) ═══
sent=re.split(r'(?<=[.]) (?=[A-Z<])',L["passage"])
pages=[sent] if NPASS==1 else [sent[:len(sent)//2], sent[len(sent)//2:]]
for pi,pgtoks in enumerate(pages):
    s=slide(); suffix=f" ({pi+1}/{NPASS})" if NPASS>1 else ""
    head(s,"READ","지문 분석 — 연결어·핵심개념·흐름"+suffix,pg())
    lx=0.55
    for lg in L.get("legend",[]):
        rect(s,lx,1.05,0.26,0.26,SW.get(lg["sw"],GOLD))
        tf=box(s,lx+0.31,1.02,3.6,0.32,MSO_ANCHOR.MIDDLE); tf.word_wrap=False; run(tf.paragraphs[0],lg["label"],11,MUTE,name=KO)
        lx+=0.72+0.14*len(lg["label"])
    pc,ptf=rect(s,0.5,1.5,12.33,5.06,PANEL,LINE,1.0)
    apsz=17 if plen(L["passage"])<1150 else (16 if plen(L["passage"])<1500 else 15)
    p=ptf.paragraphs[0]; p.line_spacing=1.32; emit(p," ".join(pgtoks),apsz,WHITE)
    if L.get("footnotes"):
        tf=box(s,0.5,6.62,12.3,0.4); run(tf.paragraphs[0],"   ".join("* "+f for f in L["footnotes"]),11,MUTE,it=True,name=KO)
    foot(s,"색: 연결어(초록)·핵심개념·구조어 — 흐름을 눈으로 잡는다")

# ═══ 직독직해 (문장별 en+ko + 구문) — 페이블 보강 슬롯 ═══
if L.get("sents"):
    sents=L["sents"]; n=len(sents)
    ppg = 1 if n<=8 else 2
    chunks=[sents] if ppg==1 else [sents[:(n+1)//2], sents[(n+1)//2:]]
    for ci,chunk in enumerate(chunks):
        s=slide(); suffix=f" ({ci+1}/{ppg})" if ppg>1 else ""
        head(s,"PARSE","직독직해 — 문장별 해석·구문"+suffix,pg())
        pc,ptf=rect(s,0.5,1.1,12.33,5.7,PANEL,LINE,1.0)
        esz = 13.5 if len(chunk)<=5 else (12 if len(chunk)<=7 else 11)
        for i,se in enumerate(chunk):
            en=se.get("en") or se.get("svo") or ""; ko=se.get("ko",""); note=se.get("note","")
            p=ptf.paragraphs[0] if i==0 else ptf.add_paragraph(); p.space_before=Pt(0 if i==0 else 6); p.line_spacing=1.18
            run(p,f"{se.get('n',i+1)}. ",esz-1,GOLD,True,name=EN)
            emit(p,en,esz,WHITE)            # en에 <<S..|tc=cyan>> 등 구문 색분해 마크업 지원
            if ko:
                p2=ptf.add_paragraph(); p2.space_before=Pt(1); p2.line_spacing=1.16
                run(p2,"   ",esz,WHITE); run(p2,strip_markup(ko),esz-0.5,TC["cyan"],name=KO)
            if note:
                p3=ptf.add_paragraph(); p3.space_before=Pt(0); p3.line_spacing=1.1
                run(p3,"   ▸ ",esz-1,GOLD); run(p3,strip_markup(note),esz-1.5,MUTE,name=KO)
        foot(s,"직독직해: S(주어)·V(동사)·O/C·M(수식) 구문을 색으로 끊어 읽는다")

# ═══ 글의 흐름 지도 ═══
if L.get("flow"):
    s=slide(); head(s,"FLOW","글의 흐름 지도",pg())
    fl=L["flow"]; n=len(fl); gap=0.12; H0=min(0.82,(5.85-(n-1)*gap)/n); y=1.05
    for f in fl:
        cl=SW.get(f.get("cl"),GOLD); cc,tf=rect(s,1.35,y,10.6,H0,PANEL,cl,2.0,anchor=MSO_ANCHOR.MIDDLE)
        p=tf.paragraphs[0]; run(p,f" ({f.get('n')}) ",13.5,cl,True,name=EN); run(p,f.get("role","")+"   ",13.5,cl,True,name=KO)
        run(p,f.get("note",""),13,WHITE,name=KO)
        if f is not fl[-1]:
            a=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(6.5),Inches(y+H0+0.002),Inches(0.22),Inches(min(0.1,gap-0.02)))
            a.fill.solid(); a.fill.fore_color.rgb=GOLD; a.line.fill.background(); a.shadow.inherit=False
        y+=H0+gap
    foot(s,"어느 문장이 주제이고 어디서 방향이 바뀌는가")

# ═══ 정답 ═══
if KIND=="reading":
    s=slide(); head(s,"ANSWER","정답 · 논리 구조",pg())
    cc,tf=rect(s,0.55,1.1,2.25,1.55,PANEL,GOLD,2.4,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run(p,"정답\n",15,MUTE,True,name=KO); run(p,CIRC[L.get("answer",1)-1],40,GOLD,True,name=EN)
    cc,tf=rect(s,3.05,1.1,9.78,1.55,PANEL2,LINE,1.0,anchor=MSO_ANCHOR.MIDDLE)
    p=tf.paragraphs[0]; p.line_spacing=1.18; run(p,L.get("answer_label","정답")+" = ",13.5,MUTE,True,name=KO); emit(p,L.get("answer_plain",""),15,WHITE)
    p2=tf.add_paragraph(); p2.space_before=Pt(3); p2.line_spacing=1.2; run(p2,L.get("rationale",""),12,GOLDL,name=KO)
    lg=L.get("logic",{})
    cc,tf=rect(s,0.55,2.82,12.28,1.78,PANEL,GOLD,1.4)
    run(tf.paragraphs[0],"논리 구조",12,GOLD,True,name=KO,spc=1.2)
    for row in lg.get("rows",[]):
        p=tf.add_paragraph(); p.space_before=Pt(5); p.line_spacing=1.18
        run(p,"  "+row.get("arrow","")+"  ",14,SW.get(row.get("acol"),GOLD),True,name=KO)
        run(p,row.get("expr","")+"  ",15,WHITE,True,name=KO)
        run(p,"["+row.get("tag","")+"]  ",12.5,TC.get(row.get("acol","gold"),GOLDL),True,name=EN)
        run(p,row.get("note",""),11,MUTE,it=True,name=EN)
    if lg.get("banner"):
        cc,tf=rect(s,0.55,4.72,12.28,0.6,GOLD,anchor=MSO_ANCHOR.MIDDLE)
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run(p,"⇅  "+lg["banner"],14.5,DARK,True,name=KO)
    cc,tf=rect(s,0.55,5.42,12.28,1.5,PANEL2,LINE,1.0)
    run(tf.paragraphs[0],"오답 근거",11.5,MUTE,True,name=KO,spc=1.0)
    for w in L.get("wrong",[]):
        p=tf.add_paragraph(); p.space_after=Pt(1); p.line_spacing=1.14
        run(p,CIRC[w.get("n",1)-1]+" ",12,TC["red"],True,name=EN); run(p,w.get("text",""),11.5,H("C7CDD6"),name=KO)
else:
    s=slide(); head(s,"ANSWER","어법 5포인트 · 정답",pg())
    pts=L.get("points",[]); n=len(pts); H0=min(0.86,4.35/max(n,1)); y=1.12
    for pt in pts:
        isX=str(pt.get("verdict","")).upper()=="X"; cl=TC["red"] if isX else TC["green"]; base=SW["red"] if isX else SW["green"]
        cc,tf=rect(s,0.55,y,0.62,H0,base,anchor=MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run(p,CIRC[pt.get("n",1)-1],16,WHITE,True,name=EN)
        cc,tf=rect(s,1.25,y,2.45,H0,PANEL2,LINE,1.0,anchor=MSO_ANCHOR.MIDDLE)
        p=tf.paragraphs[0]; run(p,pt.get("underline",""),13.5,(TC["red"] if isX else WHITE),True,ul=isX,name=EN)
        if isX and pt.get("correct"): run(p,"  → "+pt.get("correct"),12.5,TC["green"],True,name=EN)
        cc,tf=rect(s,3.78,y,2.5,H0,PANEL,cl,1.4,anchor=MSO_ANCHOR.MIDDLE); run(tf.paragraphs[0],pt.get("label",""),12.5,cl,True,name=KO)
        cc,tf=rect(s,6.36,y,5.44,H0,PANEL2,LINE,1.0,anchor=MSO_ANCHOR.MIDDLE); run(tf.paragraphs[0],pt.get("why",""),10.5,H("C7CDD6"),name=KO)
        cc,tf=rect(s,11.9,y,0.9,H0,base,anchor=MSO_ANCHOR.MIDDLE); p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER; run(p,"X" if isX else "O",17,WHITE,True,name=EN)
        y+=H0+0.075
    cc,tf=rect(s,0.55,y+0.04,12.28,6.9-(y+0.04),PANEL,GOLD,1.4)
    p=tf.paragraphs[0]; p.line_spacing=1.22; run(p,"정답 "+CIRC[L.get("answer",1)-1]+"  ",15,GOLD,True,name=EN); run(p,L.get("explanation_core",""),12.5,WHITE,name=KO)
foot(s)

prs.save(OUT); print("SAVED",OUT,"/",len(prs.slides._sldIdLst),"slides · kind=",KIND)
